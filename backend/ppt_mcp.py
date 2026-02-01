from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

mcp = FastMCP("PPT-Generator")

def set_background_color(slide, color_rgb):
    """Auxiliary function to set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

@mcp.tool()
def create_fancy_ppt(filename: str, title: str, subtitle: str = "", slides_data: list = None) -> str:
    """
    Create a PowerPoint file with a professional purple theme.
    'slides_data' is a list of dictionaries with 'title' and 'content' keys.
    """
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    
    filepath = os.path.join(os.getcwd(), filename)
    prs = Presentation()
    
    # Define Brand Colors (Deep Purple Tone)
    HEADER_PURPLE = RGBColor(45, 0, 75)   # Very deep purple
    ACCENT_PURPLE = RGBColor(100, 50, 150) # Medium purple
    BG_PURPLE = RGBColor(245, 240, 255)    # Very light purple background
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_background_color(slide, HEADER_PURPLE)
    
    title_obj = slide.shapes.title
    title_obj.text = title
    title_obj.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_obj.text_frame.paragraphs[0].font.size = Pt(44)
    title_obj.text_frame.paragraphs[0].font.bold = True

    if subtitle:
        sub = slide.placeholders[1]
        sub.text = subtitle
        sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 180, 255)
        sub.text_frame.paragraphs[0].font.size = Pt(24)

    # 2. Content Slides
    if slides_data:
        bullet_slide_layout = prs.slide_layouts[1]
        for s in slides_data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            set_background_color(slide, BG_PURPLE)
            
            # Slide Title
            title_shape = slide.shapes.title
            title_shape.text = s.get("title", "")
            title_shape.text_frame.paragraphs[0].font.color.rgb = HEADER_PURPLE
            title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Content Text
            body_shape = slide.placeholders[1]
            body_shape.text = s.get("content", "")
            for paragraph in body_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(50, 50, 50) # Dark gray for readability
                
            # Add a small purple decorative bar at the bottom
            left = Inches(0)
            top = Inches(7.2)
            width = Inches(10)
            height = Inches(0.3)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT_PURPLE
            shape.line.visible = False

    prs.save(filepath)
    return f"Professional Purple PowerPoint created: {filepath}"

if __name__ == "__main__":
    mcp.run()
